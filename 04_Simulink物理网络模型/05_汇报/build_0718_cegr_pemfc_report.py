"""Build the first two-slide Route A cEGR-PEMFC work-report deck.

The input deck supplies the corporate theme. All diagrams on slide 1 are native
PowerPoint shapes so their Chinese labels and physical flow can be edited later.
Slide 2 plots verified Stage 1 tail-window data documented in the active route.
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPORT_DIR = Path(__file__).resolve().parent
INPUT_DECK = REPORT_DIR / "0718-CEGR-PEMFC-Simulink.pptx"
OUTPUT_DECK = REPORT_DIR / "0718-CEGR-PEMFC-Simulink_v01.pptx"
ASSET_DIR = REPORT_DIR / "assets"

FONT = "Microsoft YaHei"
BLUE = "003B90"
BLUE_LIGHT = "E8F1FA"
CYAN = "007E8A"
CYAN_LIGHT = "E6F5F6"
RED = "C00000"
RED_LIGHT = "F9E8E8"
YELLOW = "D99000"
YELLOW_LIGHT = "FFF4D6"
DARK = "253342"
GRAY = "63717E"
GRAY_LIGHT = "F2F5F7"
LINE = "CBD4DC"
WHITE = "FFFFFF"
GREEN = "1E7B53"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_shape_style(shape, fill: str, line: str | None = None, width: float = 0.75):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line or fill)
    shape.line.width = Pt(width)


def set_text(
    shape,
    text: str,
    size: float,
    fill: str = DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.06,
):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin / 2)
    frame.margin_bottom = Inches(margin / 2)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)


def add_textbox(slide, x, y, w, h, text, size, fill=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box, text, size, fill, bold, align)
    return box


def add_box(slide, x, y, w, h, title, body, accent, light, title_size=15, body_size=10.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_style(shape, light, LINE, 0.6)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_shape_style(accent_bar, accent, accent, 0)
    if h < 0.65:
        set_text(shape, title + "  |  " + body, 8.5, DARK, False, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0.13)
        return shape
    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.11), Inches(w - 0.28), Inches(0.28))
    set_text(title_box, title, title_size, accent, True)
    body_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.42), Inches(w - 0.28), Inches(h - 0.5))
    set_text(body_box, body, body_size, DARK, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, 0.02)
    return shape


def add_module(slide, x, y, w, h, title, subtitle, fill, outline, text_color=DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_style(shape, fill, outline, 1.0)
    text = title if not subtitle else title + "\n" + subtitle
    set_text(shape, text, 13 if subtitle else 14, text_color, True, PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x, y, w, h, fill, direction="right"):
    mapping = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    shape = slide.shapes.add_shape(mapping[direction], Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_style(shape, fill, fill, 0.1)
    return shape


def add_line(slide, x1, y1, x2, y2, fill, dash=False, width=1.3):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color(fill)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def clear_placeholders(slide):
    for placeholder in slide.placeholders:
        if hasattr(placeholder, "text_frame"):
            placeholder.text_frame.clear()


def remove_slide(prs, index: int):
    slide_id_list = prs.slides._sldIdLst
    slide_id = list(slide_id_list)[index]
    prs.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def add_slide_header(slide, title: str, subtitle: str, page: int):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(WHITE)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.07))
    set_shape_style(top_bar, BLUE, BLUE, 0)
    add_textbox(slide, 0.48, 0.22, 10.8, 0.38, title, 24, BLUE, True)
    add_textbox(slide, 0.50, 0.68, 11.5, 0.30, subtitle, 10.5, GRAY)
    add_line(slide, 0.48, 1.05, 12.85, 1.05, LINE, False, 0.8)
    add_textbox(slide, 12.45, 0.23, 0.42, 0.28, f"0{page}", 10, GRAY, False, PP_ALIGN.RIGHT)


def add_evidence_matrix(slide, x, y, w, h):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_style(outer, "FBFCFD", "D7E1E9", 0.8)
    add_textbox(slide, x + 0.18, y + 0.10, w - 0.36, 0.23, "恒 28 A：cEGR 比矩阵已验证结果", 12.5, BLUE, True)
    add_textbox(slide, x + 0.18, y + 0.37, w - 0.36, 0.20, "同一运行态初值，600 s，尾窗 [540,600) s", 8.6, GRAY)

    row_labels = ["堆电压 (V)", "堆功率 (kW)", "入堆 O2 质量分数", "lambda_ca_in 最低值"]
    values = [["424.7", "423.9", "421.4"], ["11.892", "11.869", "11.798"], ["0.2159", "0.2060", "0.1752"], ["4.20", "4.06", "3.62"]]
    headers = ["cEGR 0", "cEGR 0.10", "cEGR 0.30"]
    left = x + 0.20
    top = y + 0.66
    label_w = 1.86
    cell_w = (w - 0.40 - label_w) / 3
    row_h = 0.37
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(label_w), Inches(row_h))
    set_shape_style(header, GRAY_LIGHT, GRAY_LIGHT, 0)
    set_text(header, "尾窗均值", 9.5, GRAY, True, PP_ALIGN.CENTER)
    for index, label in enumerate(headers):
        fill = BLUE_LIGHT if index < 2 else "E4F4F7"
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + label_w + index * cell_w), Inches(top), Inches(cell_w), Inches(row_h))
        set_shape_style(header, fill, fill, 0)
        set_text(header, label, 9.5, BLUE if index < 2 else CYAN, True, PP_ALIGN.CENTER)
    for row, label in enumerate(row_labels):
        y0 = top + row_h * (row + 1)
        label_cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(y0), Inches(label_w), Inches(row_h))
        set_shape_style(label_cell, "F6F8FA", WHITE, 0.4)
        set_text(label_cell, label, 8.8, GRAY, True, PP_ALIGN.LEFT)
        for col, value in enumerate(values[row]):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + label_w + col * cell_w), Inches(y0), Inches(cell_w), Inches(row_h))
            set_shape_style(cell, WHITE, "E4EAF0", 0.4)
            accent = BLUE if col == 0 else (CYAN if col == 1 else RED)
            set_text(cell, value, 10.5, accent, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.22, y + h - 0.30, w - 0.44, 0.20, "随 cEGR 提高，当前同相位对照下入堆氧浓度、电压与功率下降；所有案例保持有限并通过现行门。", 8.3, GRAY)


def add_slide_one(prs, layout, slide=None):
    if slide is None:
        slide = prs.slides.add_slide(layout)
    clear_placeholders(slide)
    add_slide_header(
        slide,
        "cEGR-PEMFC 系统物理网络建模平台",
        "基于官方 Gas Mixture PEMFC 案例复用，面向系统性能仿真、参数匹配与策略开发",
        1,
    )

    # System boundary.
    boundary = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.32), Inches(8.38), Inches(5.4))
    set_shape_style(boundary, "FBFCFD", "D7E1E9", 0.8)
    add_textbox(slide, 0.70, 1.42, 2.8, 0.24, "Route A 单一系统物理网络", 10.5, GRAY, True)

    # Cathode air and cEGR line.
    add_module(slide, 0.72, 2.10, 1.13, 0.65, "环境空气", "边界输入", BLUE_LIGHT, BLUE)
    add_arrow(slide, 1.91, 2.30, 0.45, 0.22, BLUE)
    add_module(slide, 2.40, 2.03, 1.42, 0.78, "供氧 / 混合", "压缩机·中冷·加湿", BLUE_LIGHT, BLUE)
    add_arrow(slide, 3.86, 2.30, 0.65, 0.22, BLUE)

    stack = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.62), Inches(2.08), Inches(1.88), Inches(2.15))
    set_shape_style(stack, DARK, DARK, 0.8)
    set_text(stack, "PEMFC\n电堆核心", 18, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 4.78, 3.52, 1.55, 0.24, "反应 · 电压 · 热", 8.8, "DDE6EE", False, PP_ALIGN.CENTER)

    add_arrow(slide, 6.55, 2.30, 0.46, 0.22, BLUE)
    add_module(slide, 7.04, 2.03, 1.36, 0.78, "阴极排气", "背压·水管理", BLUE_LIGHT, BLUE)
    add_arrow(slide, 7.57, 1.72, 0.22, 0.27, BLUE, "up")
    add_module(slide, 6.83, 1.38, 1.78, 0.44, "cEGR 回流", "阀·管路·组分回混", BLUE_LIGHT, BLUE)
    add_arrow(slide, 4.01, 1.49, 2.67, 0.20, BLUE, "left")
    add_textbox(slide, 4.42, 1.20, 2.1, 0.22, "实际回流比与氧计量比审计", 8.6, BLUE, False, PP_ALIGN.CENTER)

    # Hydrogen side.
    add_module(slide, 0.72, 4.58, 1.13, 0.65, "氢源", "压力·组分", RED_LIGHT, RED)
    add_arrow(slide, 1.91, 4.80, 0.45, 0.22, RED)
    add_module(slide, 2.40, 4.51, 1.42, 0.78, "阳极 BOP", "减压·回流·Purge", RED_LIGHT, RED)
    add_arrow(slide, 3.86, 4.80, 0.65, 0.22, RED)

    # Electrical and thermal paths.
    add_arrow(slide, 6.62, 4.19, 0.40, 0.22, YELLOW)
    add_module(slide, 7.08, 4.00, 1.32, 0.65, "电气边界", "受控负载·测量", YELLOW_LIGHT, YELLOW)
    add_arrow(slide, 5.23, 4.34, 0.22, 0.40, CYAN, "down")
    add_module(slide, 4.00, 5.07, 3.17, 0.61, "热管理 BOP", "冷却回路 · 泵 · 散热器 · 堆温控制", CYAN_LIGHT, CYAN)
    add_arrow(slide, 5.99, 4.74, 0.22, 0.28, CYAN, "down")

    # Controller on the right, with command/measurement links.
    control = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.12), Inches(2.40), Inches(3.60), Inches(2.42))
    set_shape_style(control, GRAY_LIGHT, LINE, 0.9)
    add_textbox(slide, 9.42, 2.64, 2.95, 0.33, "系统控制与观测", 16, DARK, True, PP_ALIGN.CENTER)
    add_line(slide, 9.46, 3.08, 12.38, 3.08, LINE, False, 0.6)
    add_textbox(slide, 9.45, 3.22, 2.95, 1.18, "u  负载 / 空气 / cEGR / 背压 / 热管理\nw  环境与边界条件\ny  控制可用测量\nz  守恒、状态与执行器审计", 10.5, GRAY, False)
    add_line(slide, 8.40, 2.42, 9.12, 2.78, GRAY, True, 1.2)
    add_line(slide, 8.40, 4.32, 9.12, 4.10, GRAY, True, 1.2)
    add_line(slide, 7.17, 5.36, 9.12, 4.56, GRAY, True, 1.2)
    add_line(slide, 3.15, 2.03, 9.12, 3.25, GRAY, True, 1.2)

    # Bottom attributes.
    add_box(slide, 0.45, 6.92, 3.90, 0.42, "官方物理内核", "复用官方 FuelCell / Simscape 网络", BLUE, BLUE_LIGHT, 11, 8.5)
    add_box(slide, 4.55, 6.92, 3.90, 0.42, "系统级耦合", "气体 · 水 · 热 · 电与控制共同求解", CYAN, CYAN_LIGHT, 11, 8.5)
    add_box(slide, 8.65, 6.92, 4.05, 0.42, "边界清晰", "未验证的部件效率与能力明确保留为待证", YELLOW, YELLOW_LIGHT, 11, 8.5)
    return slide


def add_slide_two(prs, layout):
    slide = prs.slides.add_slide(layout)
    clear_placeholders(slide)
    add_slide_header(
        slide,
        "Stage 1 已验证能力、研发用途与协同需求",
        "只展示已完成的系统级证据；部件效率、喘振、分离效率及完整热管理能力仍待项目数据验证",
        2,
    )

    protocol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.48), Inches(1.20), Inches(12.35), Inches(0.42))
    set_shape_style(protocol, "EEF4F9", "D7E2EB", 0.6)
    set_text(protocol, "同一 mode=1 保存运行态初值  |  恒 28 A  |  cEGR 0 / 0.10 / 0.30  |  600 s 研究时长  |  尾窗 [540,600) s", 11, BLUE, True, PP_ALIGN.CENTER)

    add_evidence_matrix(slide, 0.45, 1.78, 5.70, 2.75)
    add_textbox(slide, 0.55, 4.43, 5.42, 0.28, "数据来自已完成的 v03/mode=1 三点矩阵；未将该结果外推为压缩机或分离器能力。", 8.2, GRAY)

    add_box(slide, 6.38, 1.78, 2.00, 0.88, "实际 cEGR 跟踪", "0 / 0.10 / 0.30 三点均通过当前跟踪门。", BLUE, BLUE_LIGHT, 11.5, 9.3)
    add_box(slide, 8.60, 1.78, 2.00, 0.88, "供氧与阀边界", "lambda_ca_in 最低 3.62；阀压差为正且未饱和。", CYAN, CYAN_LIGHT, 11.5, 9.3)
    add_box(slide, 10.82, 1.78, 2.00, 0.88, "守恒与水管理", "N2/O2、O2 法拉第消耗与 WM-L1+ 气相水账本通过。", GREEN, "E8F5EF", 11.5, 9.1)
    add_box(slide, 6.38, 2.86, 6.44, 1.33, "当前能给项目研发的具体输出", "系统性能：电压、功率、氧计量、温湿压与热响应。\n设备匹配：空气/回流流量、阀压差和开度、背压与边界需求。\n策略开发：负载、空气/OER、cEGR、背压和热管理控制逻辑的联合验证。", BLUE, "F6F9FC", 12.5, 10.7)

    need = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.96), Inches(12.37), Inches(1.55))
    set_shape_style(need, "FFF8EA", "E6C87B", 0.8)
    add_textbox(slide, 0.70, 5.10, 3.55, 0.26, "进入项目研发，需要协同支持", 12.5, YELLOW, True)
    add_box(slide, 0.72, 5.46, 3.75, 0.78, "1  固定首个项目边界", "目标电堆 / 功率等级、负载谱、环境与 cEGR 研究目标。", YELLOW, "FFFDF7", 11.5, 9.6)
    add_box(slide, 4.77, 5.46, 3.75, 0.78, "2  提供关键部件与试验数据", "压缩机地图、阀/管路压降、加湿/冷却边界、堆 I-V 与气路测量。", YELLOW, "FFFDF7", 11.5, 9.3)
    add_box(slide, 8.82, 5.46, 3.75, 0.78, "3  建立联合验证工作包", "明确实验接口、数据责任人与最小验证工况，形成模型-试验闭环。", YELLOW, "FFFDF7", 11.5, 9.6)

    add_textbox(slide, 0.53, 6.74, 12.15, 0.32, "下一步：恒流多负载矩阵  →  恒功率能力验证  →  恒压控制接口设计与最小回归", 12, BLUE, True, PP_ALIGN.CENTER)
    return slide


def main():
    if not INPUT_DECK.exists():
        raise FileNotFoundError(INPUT_DECK)
    presentation = Presentation(str(INPUT_DECK))
    layout = presentation.slide_layouts[2]
    if presentation.slides:
        first_slide = presentation.slides[0]
    else:
        first_slide = presentation.slides.add_slide(layout)
    add_slide_one(presentation, layout, first_slide)
    add_slide_two(presentation, layout)
    presentation.save(str(OUTPUT_DECK))
    print(OUTPUT_DECK)


if __name__ == "__main__":
    main()
